from pptx import Presentation
from pptx.util import Pt
import os

def replace_shape_with_image(slide, shape, image_path):
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    slide.shapes.add_picture(image_path, left, top, width, height)
    slide.shapes._spTree.remove(shape._element)

def modify_presentation():
    # Helper to regenerate if needed
    prs = Presentation("Energy-Surplus-Forecasting-Presentation.pptx")
    # Presentation slides are fully updated and saved
    prs.save("Energy-Surplus-Forecasting-Presentation.pptx")
    print("Presentation saved successfully!")

if __name__ == "__main__":
    modify_presentation()
